import pytesseract
from pdf2image import convert_from_path
from bs4 import BeautifulSoup
import requests
import urllib.request
from pptx import Presentation
import glob

import fitz  # PyMuPDF

def extract_text_from_pdf(pdf_path):
    document = fitz.open(pdf_path)
    full_text = ""
    for page_num in range(len(document)):
        print(page_num)
        page = document.load_page(page_num)
        text = page.get_text()
        full_text += text
    document.close()
    return full_text

def pdf_to_text(pdf_path):
    # Convert the PDF to images
    images = convert_from_path(pdf_path, dpi=72, size=(840, 1190))

    # Extract text from each image
    texts = []
    for image in images:
        text = pytesseract.image_to_string(image)
        if len(text) > 2:
            texts.append(text)

    # Combine all the text into one string
    full_text = "\n".join(texts)

    return full_text.replace("\n\n", "\n")


def clean_html(url: str) -> str:
    head = None
    body = None
    try:
        fp = urllib.request.urlopen(url)
        content = fp.read()
        soup = BeautifulSoup(content, "lxml")
        head = soup.head.get_text(strip=True)
        body = soup.body.get_text(strip=True)
        return f"""URL: {url}\nHEADER: {head}\n\n---\n\nContent: {body}"""
    except Exception as ex:
        return f"""URL: {url}\nHEADER: {head}\n\n---\n\nContent: {body} - {ex}"""


def thread_parsing(thread_id: str):
    """
    thread id format -> chan:board:thread
    """
    chan, board, thread = thread_id.split(":")
    url = f"https://boards.4channel.org/{board}/thread/{thread}"
    soup = BeautifulSoup(requests.get(url).text, "lxml")

    post_message_tags = soup.find_all(class_="postMessage")

    for tag in post_message_tags:
        br_tags = tag.find_all("br")
        for br in br_tags:
            br.replace_with("\n")
    thread_string = ""
    # Extract and print the modified text inside each "postMessage" tag
    for tag in post_message_tags:
        thread_string += tag.get_text() + "\n"
        thread_string += "------------\n"
    return thread_string


def pptx_to_text(path):
    txt = ""
    prs = Presentation(path)

    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                txt += shape.text + "\n"
        txt += shape.text + "\n----------\n"
    return txt
